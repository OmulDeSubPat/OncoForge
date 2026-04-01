from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.config import PROJECT_ROOT


REPORTS_DIR = PROJECT_ROOT / "reports"
VIS_DIR = REPORTS_DIR / "presentation_visuals_davinci_v3"
OUTPUT_PPT = REPORTS_DIR / "OncoForge_Prezentare_DaVinci_v3.pptx"


def _add_title(slide, text: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.62))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)


def _add_picture_fit(slide, image_path: Path, left, top, width, height) -> None:
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    img_ratio = img_w / img_h
    box_ratio = width / height

    if img_ratio >= box_ratio:
        pic_w = width
        pic_h = width / img_ratio
        pic_left = left
        pic_top = top + (height - pic_h) / 2
    else:
        pic_h = height
        pic_w = height * img_ratio
        pic_left = left + (width - pic_w) / 2
        pic_top = top

    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def _add_full_image_slide(prs: Presentation, title: str, image_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)
    _add_picture_fit(
        slide=slide,
        image_path=VIS_DIR / image_name,
        left=Inches(0.42),
        top=Inches(0.95),
        width=Inches(12.5),
        height=Inches(6.25),
    )


def _add_two_image_slide(
    prs: Presentation,
    title: str,
    left_image: str,
    right_image: str,
    left_label: str,
    right_label: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)

    left_label_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.88), Inches(6.0), Inches(0.34))
    left_label_box.text_frame.text = left_label
    left_label_box.text_frame.paragraphs[0].font.size = Pt(14)
    left_label_box.text_frame.paragraphs[0].font.bold = True
    left_label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)

    right_label_box = slide.shapes.add_textbox(Inches(6.86), Inches(0.88), Inches(6.0), Inches(0.34))
    right_label_box.text_frame.text = right_label
    right_label_box.text_frame.paragraphs[0].font.size = Pt(14)
    right_label_box.text_frame.paragraphs[0].font.bold = True
    right_label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)

    _add_picture_fit(
        slide=slide,
        image_path=VIS_DIR / left_image,
        left=Inches(0.45),
        top=Inches(1.2),
        width=Inches(6.1),
        height=Inches(5.95),
    )
    _add_picture_fit(
        slide=slide,
        image_path=VIS_DIR / right_image,
        left=Inches(6.78),
        top=Inches(1.2),
        width=Inches(6.1),
        height=Inches(5.95),
    )


def _add_intro_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "OncoForge - Prioritizare AI pentru inhibitori EGFR")

    subtitle = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.65))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.clear()
    p = subtitle_frame.paragraphs[0]
    p.text = "Prezentare competitie (structura 1-10)"
    p.font.size = Pt(22)
    p.font.bold = False
    p.font.color.rgb = RGBColor(51, 65, 85)

    details = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.6), Inches(3.6))
    details_frame = details.text_frame
    details_frame.clear()
    details_lines = [
        "Autori: [completeaza numele]",
        "Coordonator: [completeaza profesorul]",
        "Institutie: [completeaza institutia]",
        "Data: [completeaza data prezentarii]",
    ]
    for index, line in enumerate(details_lines):
        paragraph = details_frame.paragraphs[0] if index == 0 else details_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(30, 41, 59)

    note = slide.shapes.add_textbox(Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.5))
    note.text_frame.text = "Slide 1 ramane editabil pentru completare manuala."
    note.text_frame.paragraphs[0].font.size = Pt(13)
    note.text_frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)


def _add_future_plans_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Planuri de viitor")

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.9), Inches(5.6))
    frame = box.text_frame
    frame.clear()
    lines = [
        "Completeaza manual planurile de viitor:",
        "1. [Plan 1]",
        "2. [Plan 2]",
        "3. [Plan 3]",
        "4. [Plan 4]",
    ]
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(24 if idx == 0 else 22)
        paragraph.font.bold = idx == 0
        paragraph.font.color.rgb = RGBColor(30, 41, 59)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_intro_slide(prs)
    _add_full_image_slide(prs, "Problema abordata", "02_problema_abordata_burden.png")
    _add_two_image_slide(
        prs,
        "Context stiintific",
        "03a_tabel_medicamente.png",
        "03b_trend_publicatii.png",
        "Medicamente EGFR pe piata",
        "Trend cercetare: generare molecule",
    )
    _add_full_image_slide(prs, "Abstract: restrangerea spatiului chimic", "04_abstract_spatiu_chimic.png")
    _add_full_image_slide(prs, "Metodologia proiectului", "05_metodologie.png")
    _add_two_image_slide(
        prs,
        "Rezultate cruciale pe modele",
        "06a_rmse_modele.png",
        "06b_r2_modele.png",
        "Comparatie RMSE",
        "Comparatie R^2",
    )
    _add_full_image_slide(prs, "Rezultate multi-agent", "07_rezultate_multi_agent.png")
    _add_full_image_slide(prs, "AI-ul nostru vs alte studii", "08_ai_vs_studii.png")
    _add_future_plans_slide(prs)
    _add_full_image_slide(prs, "Referinte", "10_referinte.png")

    prs.save(OUTPUT_PPT)
    return OUTPUT_PPT


def main() -> None:
    output = build_presentation()
    print(f"[OK] PPT generated: {output}")


if __name__ == "__main__":
    main()
